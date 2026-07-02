"""Generate the 15-minute demo PowerPoint for the
IT / Ticket Auto-Resolution Agent (CP06) capstone.

Run:
    python tools/generate_ppt.py
Produces:
    IT_Ticket_Auto_Resolution_Agent_Demo.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2C, 0x4A)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
TEAL = RGBColor(0x17, 0xA2, 0x8C)
AMBER = RGBColor(0xF0, 0xA5, 0x00)
GREY = RGBColor(0x5B, 0x66, 0x70)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x33)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    return shp


def set_text(shape, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def add_textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    set_text(tb, text, size, color, bold, align, anchor, font)
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=DARK, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        run = p.add_run()
        bullet = "•  " if lvl == 0 else "–  "
        run.text = bullet + txt
        run.font.size = Pt(size - lvl * 2)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb


def header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.06), TEAL)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.55),
                title, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        add_textbox(slide, Inches(0.62), Inches(0.78), Inches(11.5), Inches(0.3),
                    kicker, size=13, color=RGBColor(0xB9, 0xD3, 0xEA))
    # page number footer
    add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
                f"{len(prs.slides._sldIdLst)}", size=10, color=GREY,
                align=PP_ALIGN.RIGHT)


def add_arrow(slide, x, y, w, h, color=GREY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_down_arrow(slide, x, y, w, h, color=GREY):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ============================================================================
# Slide 1 — Title
# ============================================================================
s = add_slide()
fill_bg(s, NAVY)
add_rect(s, 0, Inches(4.9), SLIDE_W, Inches(0.08), TEAL)
add_rect(s, 0, Inches(6.4), SLIDE_W, Inches(1.1), RGBColor(0x0A, 0x20, 0x38))
add_textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.4),
            "IT / Ticket Auto-Resolution Agent", size=48, color=WHITE, bold=True)
add_textbox(s, Inches(0.85), Inches(3.15), Inches(11.5), Inches(0.7),
            "An Agentic AI system that analyzes IT tickets, finds root cause, and "
            "recommends or auto-executes resolutions.", size=20,
            color=RGBColor(0xCF, 0xE1, 0xF0))
add_textbox(s, Inches(0.85), Inches(4.15), Inches(11), Inches(0.5),
            "Capstone CP06  •  DAAI Agentic AI — Advance CoP", size=16, color=TEAL, bold=True)
add_textbox(s, Inches(0.85), Inches(6.6), Inches(11), Inches(0.6),
            "Stack: Python  •  LangChain (sequential agents)  •  llama3.2 (Ollama)  •  RAG (ChromaDB)  •  ServiceNow REST",
            size=13, color=RGBColor(0xB9, 0xD3, 0xEA))

# ============================================================================
# Slide 2 — Problem & Objective
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "The Problem & Our Objective", "Why an agent for IT ticket resolution?")
# Two columns
add_round(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), LIGHT)
add_textbox(s, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.5),
            "The Problem", size=22, color=NAVY, bold=True)
add_bullets(s, Inches(0.9), Inches(2.4), Inches(5.3), Inches(4.0), [
    ("IT teams handle high volumes of repetitive tickets", 0),
    ("Manual triage is slow and inconsistent", 0),
    ("Knowledge is buried in runbooks & past incidents", 0),
    ("Delays increase resolution time & cost", 0),
    ("Agents spend time on low-value, known fixes", 0),
], size=17, gap=10)

add_round(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0xE9, 0xF4, 0xF1))
add_textbox(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.5),
            "Our Objective", size=22, color=TEAL, bold=True)
add_bullets(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(4.0), [
    ("Autonomously analyze & classify incoming tickets", 0),
    ("Retrieve similar past resolutions using RAG", 0),
    ("Recommend fixes with a confidence score", 0),
    ("Safely auto-resolve pre-approved ticket types", 0),
    ("Log every decision for audit & explainability", 0),
], size=17, gap=10)

# ============================================================================
# Slide 3 — Architecture (6-agent pipeline)
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "Solution Architecture", "Six specialized agents in a sequential LangChain pipeline")

stages = [
    ("1  Ingestion", "Fetch ticket\n(ServiceNow / mock)", BLUE),
    ("2  Classification", "LLM: type,\nseverity, confidence", BLUE),
    ("3  Retrieval (RAG)", "Top-k similar\nresolutions", TEAL),
    ("4  Decision", "Recommend steps\n+ confidence", TEAL),
    ("5  Execution", "Auto-resolve\nor recommend", AMBER),
]
x = Inches(0.45)
box_w = Inches(2.28)
box_h = Inches(1.7)
y = Inches(2.4)
gap = Inches(0.15)
for i, (title, desc, color) in enumerate(stages):
    b = add_round(s, x, y, box_w, box_h, color)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(11); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    if i < len(stages) - 1:
        add_arrow(s, x + box_w + Emu(5000), y + Inches(0.62), gap - Emu(10000), Inches(0.45), GREY)
    x = x + box_w + gap

# Supporting layer
add_round(s, Inches(0.45), Inches(4.5), Inches(5.9), Inches(1.0), RGBColor(0xED, 0xE7, 0xF6))
add_textbox(s, Inches(0.65), Inches(4.62), Inches(5.5), Inches(0.8),
            "Knowledge Base (ChromaDB)\nRunbooks + past incidents feed the RAG agent",
            size=13, color=DARK)
add_round(s, Inches(6.6), Inches(4.5), Inches(6.1), Inches(1.0), RGBColor(0xFC, 0xED, 0xD5))
add_textbox(s, Inches(6.8), Inches(4.62), Inches(5.7), Inches(0.8),
            "llama3.2 via Ollama (free/unlimited)\nNetwork-first, with automatic offline fallback",
            size=13, color=DARK)

add_round(s, Inches(0.45), Inches(5.7), Inches(12.25), Inches(0.9), NAVY)
add_textbox(s, Inches(0.65), Inches(5.82), Inches(11.9), Inches(0.7),
            "Logging & Monitoring  —  every stage writes a structured JSON audit entry (explainability by design)",
            size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# Slide 4 — Decision logic (auto-resolve vs recommend)
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "Decision Logic: Auto-Resolve vs. Escalate", "Safe automation with a confidence gate")

# center decision diamond
add_textbox(s, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
            "Ticket flows in →", size=16, color=GREY, bold=True)
dia = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4.9), Inches(1.9), Inches(3.5), Inches(1.9))
dia.fill.solid(); dia.fill.fore_color.rgb = AMBER; dia.line.fill.background(); dia.shadow.inherit = False
tf = dia.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Whitelisted type\nAND confidence ≥ threshold?"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"

# Yes branch
add_round(s, Inches(1.0), Inches(4.4), Inches(5.0), Inches(1.4), TEAL)
add_textbox(s, Inches(1.2), Inches(4.55), Inches(4.6), Inches(1.1),
            "YES → Execution Agent\nApply safe automated fix, set ticket to Resolved, comment steps taken",
            size=14, color=WHITE, bold=True)
# No branch
add_round(s, Inches(7.3), Inches(4.4), Inches(5.0), Inches(1.4), BLUE)
add_textbox(s, Inches(7.5), Inches(4.55), Inches(4.6), Inches(1.1),
            "NO → Recommend & Escalate\nPost recommended steps, assign to human agent for approval",
            size=14, color=WHITE, bold=True)

add_down_arrow(s, Inches(3.2), Inches(3.9), Inches(0.5), Inches(0.45), TEAL)
add_down_arrow(s, Inches(9.6), Inches(3.9), Inches(0.5), Inches(0.45), BLUE)

add_round(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.8), LIGHT)
add_textbox(s, Inches(1.2), Inches(6.2), Inches(11), Inches(0.6),
            "Auto-resolve whitelist: Password Reset • VPN Connectivity • Disk Space Cleanup    "
            "|    Always human: Access/Permission Requests",
            size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# Slide 5 — Tech stack
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "Technology Stack", "All open-source or provided — no paid licenses")
rows = [
    ("LLM", "llama3.2 via Ollama (provided endpoint)", "Free tokens; offline fallback"),
    ("Orchestration", "LangChain sequential agents", "Simple, explainable pipeline"),
    ("RAG", "ChromaDB + Ollama / TF-IDF embeddings", "Network-first, offline-safe"),
    ("ServiceNow", "REST client stub + mock JSON store", "Demo without live credentials"),
    ("Config & Secrets", "YAML + .env", "Secure credential handling"),
    ("Testing", "14 unittest tests (incl. guardrails)", "Validates safe automation"),
]
y = Inches(1.55)
# header row
add_rect(s, Inches(0.6), y, Inches(3.0), Inches(0.55), NAVY)
add_rect(s, Inches(3.6), y, Inches(4.6), Inches(0.55), NAVY)
add_rect(s, Inches(8.2), y, Inches(4.5), Inches(0.55), NAVY)
for txt, xx, ww in [("Layer", Inches(0.6), Inches(3.0)), ("Tool / Approach", Inches(3.6), Inches(4.6)), ("Why", Inches(8.2), Inches(4.5))]:
    add_textbox(s, xx + Inches(0.15), y, ww, Inches(0.55), txt, size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
y = y + Inches(0.55)
for i, (a, b, c) in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), y, Inches(3.0), Inches(0.78), bg)
    add_rect(s, Inches(3.6), y, Inches(4.6), Inches(0.78), bg)
    add_rect(s, Inches(8.2), y, Inches(4.5), Inches(0.78), bg)
    add_textbox(s, Inches(0.75), y, Inches(2.8), Inches(0.78), a, size=14, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(3.75), y, Inches(4.4), Inches(0.78), b, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(8.35), y, Inches(4.3), Inches(0.78), c, size=13, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.78)

# ============================================================================
# Slide 6 — Live Demo flow
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "Live Demo", "Two tickets, end-to-end")

# Demo 1
add_round(s, Inches(0.6), Inches(1.55), Inches(5.85), Inches(4.9), RGBColor(0xE9, 0xF4, 0xF1))
add_textbox(s, Inches(0.85), Inches(1.7), Inches(5.4), Inches(0.5), "Demo 1 — Auto-Resolve", size=20, color=TEAL, bold=True)
add_bullets(s, Inches(0.9), Inches(2.4), Inches(5.3), Inches(4.0), [
    ("Ingest a 'Password Reset' ticket", 0),
    ("LLM classifies: Account, High confidence", 0),
    ("RAG retrieves the reset runbook", 0),
    ("Decision: whitelisted + confident", 0),
    ("Execution auto-resolves the ticket", 0),
    ("Status → Resolved, comment added", 0),
], size=16, gap=9)

# Demo 2
add_round(s, Inches(6.9), Inches(1.55), Inches(5.85), Inches(4.9), LIGHT)
add_textbox(s, Inches(7.15), Inches(1.7), Inches(5.4), Inches(0.5), "Demo 2 — Recommend & Escalate", size=20, color=BLUE, bold=True)
add_bullets(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(4.0), [
    ("Ingest an 'Access Request' ticket", 0),
    ("LLM classifies: Access / Permission", 0),
    ("RAG retrieves the approval runbook", 0),
    ("Decision: NOT auto-resolvable", 0),
    ("Recommended steps posted to ticket", 0),
    ("Assigned to human for approval", 0),
], size=16, gap=9)

# ============================================================================
# Slide 7 — Explainability / audit
# ============================================================================
s = add_slide(); fill_bg(s, WHITE)
header(s, "Explainability & Audit Trail", "Every decision is traceable")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(5.6), Inches(4.8), [
    ("Structured JSON log per ticket", 0),
    ("Captures classification + confidence", 0),
    ("Records retrieved knowledge sources", 0),
    ("Logs the auto-resolve / escalate reason", 0),
    ("Timestamps every agent action", 0),
    ("Enables audit, debugging & trust", 0),
], size=17, gap=11)
# code-ish panel
add_round(s, Inches(6.6), Inches(1.6), Inches(6.1), Inches(4.9), DARK)
code = (
    "{\n"
    '  "ticket_id": "INC0012",\n'
    '  "category": "Password Reset",\n'
    '  "severity": "High",\n'
    '  "confidence": 0.94,\n'
    '  "retrieved": ["kb/password_reset.md"],\n'
    '  "decision": "auto_resolve",\n'
    '  "action": "reset_password",\n'
    '  "status": "Resolved",\n'
    '  "timestamp": "2026-07-01T10:22Z"\n'
    "}"
)
tb = s.shapes.add_textbox(Inches(6.85), Inches(1.8), Inches(5.6), Inches(4.5))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(code.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.name = "Consolas"; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x9C, 0xDC, 0xFE)

# ============================================================================
# Slide 8 — Roadmap / learnings / thank you
# ============================================================================
s = add_slide(); fill_bg(s, NAVY)
add_rect(s, 0, Inches(1.15), SLIDE_W, Inches(0.06), TEAL)
add_textbox(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.6),
            "Learnings & Next Steps", size=30, color=WHITE, bold=True)
add_bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.6), [
    ("Agentic pipelines make IT automation explainable and safe", 0),
    ("RAG grounds recommendations in real runbooks — fewer hallucinations", 0),
    ("Confidence + whitelist gate keeps auto-resolution safe", 0),
    ("Network-first design falls back to offline mode with no code changes", 0),
    ("14 automated tests validate the safety guardrails", 0),
    ("Next: connect the real ServiceNow instance via the REST client", 0),
    ("Next: expand knowledge base & add more auto-resolve categories", 0),
], size=18, color=RGBColor(0xDC, 0xE9, 0xF4), gap=12)
add_rect(s, 0, Inches(6.3), SLIDE_W, Inches(1.2), RGBColor(0x0A, 0x20, 0x38))
add_textbox(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7),
            "Thank you  —  IT / Ticket Auto-Resolution Agent  •  Questions?",
            size=20, color=TEAL, bold=True)

# ----------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "IT_Ticket_Auto_Resolution_Agent_Demo.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
